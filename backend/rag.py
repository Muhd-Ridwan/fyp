"""
RAG pipeline for FYP
Handles:
    1. Text extraction per file type
    2. Chunking
    3. Embedding via Bedrock Titan V2
    4. Upserting vectors into Pinecone with department -> file metadata
    5. Querying Pinecone for relevant chunks
"""

import io
import os
import pdfplumber
import openpyxl
import pytesseract
import csv
import logging

from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
from pinecone import Pinecone
from pinecone.errors.exceptions import NotFoundError

from config import PINECONE_API_KEY, PINECONE_INDEX_NAME
from bedrock_client import get_embedding

logger = logging.getLogger(__name__)

# Tesseract binary path - Window local dev only
# On EB (Amazon Linux), tesseract is on PATH after yum install, no path needed
if os.name == "nt":
    pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

_pc = Pinecone(api_key=PINECONE_API_KEY)
_index = _pc.Index(PINECONE_INDEX_NAME)

CHUNK_SIZE = 500
CHUNK_OVERLAP = 50


def _extract_text(file_bytes: bytes, display_name: str) -> str:
    ext = display_name.rsplit(".", 1)[-1].lower() if "." in display_name else ""

    if ext == "pdf":
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)
    
    if ext == "docx":
        doc = DocxDocument(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs)
    
    if ext == "xlsx":
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text:
                    lines.append(row_text)
        return "\n".join(lines)
    
    if ext == "pptx":
        prs = Presentation(io.BytesIO(file_bytes))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    lines.append(shape.text)
        return "\n".join(lines)
    
    if ext in ("png", "jpg", "jpeg"):
        image = Image.open(io.BytesIO(file_bytes))
        return pytesseract.image_to_string(image)
    
    if ext == "csv":
        text = file_bytes.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(text))
        return "\n".join(" | ".join(row) for row in reader)
    
    return file_bytes.decode("utf-8", errors="replace")

def _chunk_text(text: str) -> list[str]:
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + CHUNK_SIZE, len(words))
        chunks.append(" ".join(words[start:end]))
        start += CHUNK_SIZE - CHUNK_OVERLAP
    return [c for c in chunks if c.strip()]

def index_document(file_bytes: bytes, display_name: str, file_id: str, department: str) -> None:
    text = _extract_text(file_bytes, display_name)
    if not text.strip():
        return
    
    chunks = _chunk_text(text)
    vectors = []
    for i, chunk in enumerate(chunks):
        embedding = get_embedding(chunk)
        vectors.append({
            "id": f"{file_id}_{i}",
            "values": embedding,
            "metadata": {
                "department": department,
                "file_id": file_id,
                "display_name": display_name,
                "chunk_index": i,
                "text": chunk,
            },
        })
    for batch_start in range(0, len(vectors), 100):
        _index.upsert(vectors=vectors[batch_start:batch_start + 100])

def delete_document_vectors(file_id: str, department: str) -> None:
    try:
        _index.delete(filter={"file_id": {"$eq": file_id}, "department": {"$eq": department}})
    except NotFoundError:
        logger.warning("No vectors found for file %s in dept %s - skipping", file_id, department)

def query_documents(question: str, department: str, top_k: int = 10) -> list[dict]:
    embedding = get_embedding(question)
    result = _index.query(
        vector=embedding,
        top_k=top_k,
        filter={"department": {"$eq": department}},
        include_metadata=True,
    )
    return [
        {
            "text": match["metadata"]["text"],
            "display_name": match["metadata"].get("display_name", "Unknown"),
            "filename": match["metadata"].get("display_name", "Unknown"),
        }
        for match in result.get("matches", [])
        if "text" in match.get("metadata", {})
    ]